#-*- coding:utf-8 -*-
from pptx import Presentation
from pptx.util import Pt
import json
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
import requests
from server.awsModule import *
from gpuEngine.ner_api import *
import random
import os

class PPTData:
    def __init__(self, slideList,pptTopic,toc):
        self._slideList = slideList
        # initial value
        self._topic = pptTopic
        self._slideTypes = slideList
        self._toc = toc

        # 토픽에 맞는 피피티 템플릿 다운로드
        try :
            # baseFilePath = "basePPT/"+ self._topic + "_"+str(random.randint(1,2))+".pptx"
            baseFilePath = "basePPT/"+ self._topic + "_1.pptx"
            # self._resFilePath = "pptEngine/"+ self._topic + "_"+str(random.randint(1,2))+".pptx"
            self._resFilePath = "pptEngine/"+ self._topic + "_1.pptx"
            downloadFileFromS3(baseFilePath, self._resFilePath)
            self._basePrs = Presentation(self._resFilePath)
        except :
            baseFilePath = "basePPT/DEFAULT_1.pptx"
            self._resFilePath = "pptEngine/DEFAULT_1.pptx"
            downloadFileFromS3(baseFilePath, self._resFilePath)
            self._basePrs = Presentation(self._resFilePath)

        # downloadFileFromS3(baseFilePath, self._resFilePath)
        # self._basePrs = Presentation(self._resFilePath)

    # 전체 주제 getter
    @property
    def topic_(self):
        return self._topic

    # 전체 주제 setter
    #@topic.setter
    def topic(self):
        # GPU 서버로 textData를 보내서 NLP 결과를 받아온다
        result = None
        # 결과를 토대로 topic을 선정한다
        self._topic = result
            
    def newSlide(self, slideType):
        slide = self._basePrs.slides.add_slide(self._basePrs.slide_layouts[slideType])  # ppt 객체, 슬라이드마스터 번호, 제목
        #slide.shapes.title.text = self._textData._mainTitle
        return slide

    # 제목 슬라이드를 만드는 함수
    def titleSlide(self):  # NewSlide, 제목, 부제목
        slide = self.newSlide(0)#2에 제목 넣자
        slide.shapes.placeholders[0].text = self._textData._mainTitle
        slide.shapes.placeholders[1].text = self._textData._subTitle
        return slide


    # 슬라이드 객체의 특정 텍스트박스에 대한 설정
    def setTextBox(self, slide, cnt, fontAdr):  # 슬라이드 객체, 텍스트박스 번호, 폰트주소
        text_box = slide.shapes.placeholders[cnt].text_frame
        text_box.word_wrap = True
        text_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        #p = slide.shapes.add_textbox(1,1,1,1)
        #text_box.fit_text(font_file=fontAdr)
        return text_box

    # 텍스트박스 객체에 새로운 텍스트를 추가한다. (줄단위)

    def firstLine(self, textbox, text, font, size):
        #textbox.font.name = font
        #textbox.font.size = Pt(size)
        if isinstance(text, list):
            for t in text:
                textbox.text += '\n' if t!=text[0] else ''
                textbox.text = textbox.text + t
        else:
            textbox.text = text
        #textbox.alignment = PP_ALIGN.LEFT
        return textbox

    def newLine(self, textbox, text, font, size):  # 텍스트박스 객체, 내용, 폰트, 크기
        line = textbox.add_paragraph()
        #line.font.name = font
        #line.font.size = Pt(size)
        line.text = text
        #line.alignment = PP_ALIGN.LEFT
        return textbox

    def newLine_beauty(self,textbox,text,font,size,bold,rgb,center):
        line = textbox.add_paragraph()
        line.font.name = font
        line.font.size = Pt(size)
        line.text = text
        line.font.color.rgb = RGBColor(rgb[0],rgb[1],rgb[2])
        line.font.bold = bold
        if(center):
            line.alignment = PP_ALIGN.CENTER

    def transitionSlide(self,idx):
        slide = self.newSlide(2)
        text_box = self.setTextBox(slide,0,'pptEngine/static/SangSangTitleM.ttf')
        self.firstLine(text_box,self._textData._midTitles[idx],'SangSangTitleM',44)

    def index(self):
        slide = self.newSlide(1)
        text_box = self.setTextBox(slide,0,'pptEngine/static/SangSangTitleM.ttf')
        self.firstLine(text_box,'Contents','SangSangTitleM',32)
        text_box = self.setTextBox(slide, 13, 'pptEngine/static/a타이틀고딕3.ttf')
        pre_idx=0
        self.firstLine(text_box,self._textData._midTitles[0],'a타이틀고딕3',24)
        for sld_title in self._textData._slideTitles:
            if(sld_title[1]>pre_idx):
                pre_idx=sld_title[1]
                self.newLine(text_box,self._textData._midTitles[pre_idx],'a타이틀고딕3',24)
            self.newLine_beauty(text_box, '- '+sld_title[0],'a타이틀고딕3', 18,False,[0xD0,0xFF,0x80],False)

    # ppt를 저장한다.
    def write(self, file_name):
        self._basePrs.save(file_name)

    def input_title(self,slide, title):
        text_box = self.setTextBox(slide, 0, 'pptEngine/static/SangSangTitleM.ttf')
        self.firstLine(text_box, title, 'SangSangTitleM', 32)

    def timeMultiLine(self,slide_,title,headers,bodies,Links):
        # print(slide_)
        slide = self.newSlide(slide_)
        self.input_title(slide,title)
        line_cnt = 13
        for header in headers:
            text_box = self.setTextBox(slide, line_cnt,'pptEngine/static/DOSSaemmul.ttf')
            self.firstLine(text_box,header,'Arial',10)
            line_cnt = line_cnt+2
        line_cnt = 14
        for body in bodies:
            text_box = self.setTextBox(slide, line_cnt, 'pptEngine/static/DOSSaemmul.ttf')
            # print(body)
            #if flipflop ==0 :
            #    self.firstLine(text_box,body,'Arial',10)
            #    flipflop=1
            #else :
            #    self.newLine(text_box,body,'Arial',10)
            self.firstLine(text_box,body,'Arial',10)
            line_cnt = line_cnt + 2
        if not Links:
            return slide
        # print(str(slide_)+'번째 슬라이드')
        image_start = int(17+(((slide_%6)-1)/2)*2)#7,9,11 now..?-> 17,19,21
        # print(image_start)
        num_list=[]
        for i in Links:
            num_list.append(image_start)
            image_start=image_start+1
            # print(image_start)
        return self.input_image(slide,num_list,Links)

    def defaultLine(self, title, contents, Links):
        slide = self.newSlide(3)
        self.input_title(slide, title)
        text_box = self.setTextBox(slide, 1, 'pptEngine/static/a타이틀고딕3.ttf')
        flipflop=0
        for line in contents:
            if flipflop == 0 :
                self.firstLine(text_box, line[1], 'a타이틀고딕3', 20)
                flipflop = 1
            else :
                self.newLine(text_box, line[1], 'a타이틀고딕3', 20)
                
        if not Links:
            return slide
        image_start = 10
        num_list=[]
        for i in Links:
            num_list.append(image_start)
            image_start=image_start+1
        return self.input_image(slide,num_list,Links)

    def input_image(self,slide,num_ls,Links):
        n=0
        for i in num_ls:
            image_placeholder = slide.placeholders[i]
            r = requests.get(Links[n])
            n = n+1
            file = open("tempoimage.jpg","wb")
            file.write(r.content)
            file.close()
            image_placeholder.insert_picture('tempoimage.jpg')
        return slide

    def image_from_web(self,slide,blank_num,num_ls,keyword,Links):
        Links.append(wrap_image(keyword))
        num_ls.append(blank_num)
        return self.input_image(slide,num_ls,Links)

    def singleLine(self,title,lines,Links):
        # print("--In SingleLine-- with"+title)
        if not Links:
            slide = self.newSlide(4)
        else :
            slide = self.newSlide(5)
        self.input_title(slide, title)
        text_box = self.setTextBox(slide, 13, 'pptEngine/static/a타이틀고딕3.ttf')
        line = lines
        self.firstLine(text_box, line, 'a타이틀고딕3', 24)
        num_list = []
        if not Links:
            #print(line)
            #keyword = get_NNG(line)
            #print(keyword)
            #return self.image_from_web(slide, 14, num_list, keyword, Links)
            return slide
        image_start = 14
        for i in Links:
            num_list.append(image_start)
            image_start=image_start+1
        return self.input_image(slide,num_list,Links)

    def generate_slide(self, slideObj):
        print(slideObj)
        if (isinstance(slideObj,SlideType_head_default)):
            contents=[]
            for tuple in slideObj._headTuples:
                contents.append(tuple[0])
                contents.append(tuple[1])
            return self.defaultLine(slideObj._title,contents,slideObj._imageLinks)
        elif (isinstance(slideObj, SlideType_head_timeLine)):
            headers=[]
            bodies=[]
            for tuple in slideObj._headTuples:
                headers.append(tuple[0])
                body = []
                for line in tuple[1]:
                    body.append(line[1])
                bodies.append(body)
            if not slideObj._imageLinks :
                slide = 2 + len(headers) * 2#self.newSlide(2 + len(bodies) * 2)  # 6,8,10
            else :
                slide = 3+len(headers)*2 #self.newSlide(3+len(bodies)*2) # 7,9,11
            print('슬라이드 번호 : '+str(slide))
            return self.timeMultiLine(slide,slideObj._title,headers,bodies,slideObj._imageLinks)
        elif (isinstance(slideObj, SlideType_head_multiLine)):
            headers = []
            bodies = []
            for headTuple in slideObj._headTuples:
                # print(headTuple)
                headers.append(headTuple[0])
                tmp_body=[]
                for body in headTuple[1]:
                    tmp_body.append(body[1])
                bodies.append(tmp_body)
                # print(bodies)
            if not slideObj._imageLinks :
                slide = 8+len(bodies)*2#self.newSlide(8 + len(bodies) * 2)  # 6,8,10
            else :
                slide = 9+len(bodies)*2#self.newSlide(9+len(bodies)*2) # 7,9,11
            return self.timeMultiLine(slide,slideObj._title,headers,bodies,slideObj._imageLinks)
        elif (isinstance(slideObj, SlideType_head_timeLine)):
            headers = []
            bodies = []
            for tuple in slideObj._headTuples:
                headers.append(tuple[0])
                tmp_body = []
                for body in tuple[1]:
                    tmp_body.append(body[1])
                bodies.append(tmp_body)
            if not slideObj._imageLinks :
                slide = 2 + len(bodies) * 2 #self.newSlide(2 + len(bodies) * 2)  # 6,8,10
            else :
                slide = 3 + len(bodies) * 2 #self.newSlide(3+len(bodies)*2) # 7,9,11
            return self.timeMultiLine(slide,slideObj._title,headers,bodies,slideObj._imageLinks)

        elif (isinstance(slideObj,SlideType_default)):
            if not slideObj._contentsList :
                slide = self.newSlide(40)
                self.input_title(slide, slideObj._title)
                image_start = 10
                num_list = []
                num_list.append(image_start)
                return self.input_image(slide,num_list,slideObj._imageLinks)

            return self.defaultLine(slideObj._title,slideObj._contentsList,slideObj._imageLinks)

        elif (isinstance(slideObj,SlideType_singleLine)):
            return self.singleLine(slideObj._title,slideObj._text,slideObj._imageLinks)

        elif (isinstance(slideObj, SlideType_multiLine)):
            headers = []
            bodies = []
            for line in slideObj._textList:
                if not line:
                    continue
                bodies.append([line])
                print(line)
            print(str(len(bodies)) + '*2+8or9')
            if not slideObj._imageLinks :
                slide = 28 + len(bodies) * 2 #self.newSlide(8 + len(bodies) * 2)  # 6,8,10
            else :
                slide = 29+len(bodies)*2 #self.newSlide(9+len(bodies)*2) # 7,9,11
            return self.timeMultiLine(slide,slideObj._title,headers,bodies,slideObj._imageLinks)

        elif (isinstance(slideObj, SlideType_timeLine)):
           headers = []
            bodies = []
            for line in slideObj._textList:
                bodies.append([line])
            if not slideObj._imageLinks :
                slide = 22 + len(bodies) * 2#self.newSlide(2 + len(bodies) * 2)  # 6,8,10
            else :
                slide = 23+len(bodies)*2#self.newSlide(3+len(bodies)*2) # 7,9,11
            return self.timeMultiLine(slide,slideObj._title,headers,bodies,slideObj._imageLinks)

        elif (isinstance(slideObj, SlideType_title)):
            slide = self.newSlide(0)
            title = slideObj._titleTuple[0]
            subTitle = slideObj._titleTuple[1]
            text_box = self.setTextBox(slide, 0, 'pptEngine/static/SangSangTitleM.ttf')
            self.firstLine(text_box, title, 'SangSangTitleM', 32)
            text_box = self.setTextBox(slide, 1, 'pptEngine/static/SangSangTitleM.ttf')
            self.firstLine(text_box, subTitle, 'SangSangTitleM', 32)
            #if slideObj._imageLinks: TODO

        elif (isinstance(slideObj,SlideType_midTitle)):
            slide = self.newSlide(2)#간지(?)
            midTitle = slideObj._midTitle
            text_box = self.setTextBox(slide, 0, 'pptEngine/static/SangSangTitleM.ttf')
            self.firstLine(text_box, midTitle, 'SangSangTitleM', 32)
            #if slideObj._imageLinks: TODO

        elif (isinstance(slideObj,SlideType_title)):
            slide = self.newSlide(0)  # 2에 제목 넣자
            slide.shapes.placeholders[0].text = slideObj._titleTuple[0]
            slide.shapes.placeholders[1].text = slideObj._titleTuple[1]

    def generate(self):
        # self.basePrs()
        #self.titleSlide()
        #self.index()
        slide_idx = -1
        mid_slide = -1
        for slide_ in self._slideTypes:
            slide_idx = slide_idx + 1
#            if(self._textData._slideTitles[slide_idx][1]>mid_slide):
#                mid_slide = mid_slide + 1
#                self.transitionSlide(mid_slide)
            self.generate_slide(slide_)
        os.remove(self._resFilePath)

    def idx_check(self,idx):
        slide = self.newSlide(idx)
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # print('%d,%s'%(phf.idx,phf.type))


def wrap_image(keywords):
    API_KEY = '9550200-7709f1b4e4c3d4b4c800b8188'
    #gmail_KEY
    image = Image(API_KEY)
    ims = image.search(q=keywords,
                       lang='en',
                       image_type='photo',
                       orientation='horizontal',
                       category='all',
                       safesearch='true',
                       order='popular',
                       page=1,
                       per_page=3)
    if(ims['total']<4):
        return 'https://cdn.pixabay.com/photo/2020/06/30/23/21/cat-5357876_150.jpg'
    tmp = ims['hits'][0]
    url = tmp['previewURL']
    return url


# 디폴트 타입
# [String] - 한줄내용
class SlideType:
    def __init__(self, lines):
        self._lines = lines
    def generate(self):
        print('generated')

# # 타임라인 타입 (일정, 과정, 단계)
# # [(String, String)] - (시간, 한줄내용)
# class SlideType_timeline(SlideType):
#     def __init__(self, timeTuples):
#         self._timeTuples = timeTuples


# # h5 타입 (비교대조 등)
# # [(String,[String])] - (헤딩,[내용들])
# class SlideType_h5(SlideType):
#     def __init__(self, h5Tuples):
#         self._h5Tuples = h5Tuples
#         self._lines = h5Tuples[0][0]


# # 정의 타입 (? -> "")
# # String
# class SlideType_definition(SlideType):
#     def __init__(self, defStr):
#         self._defStr = defStr

# # def slide_test():
# #     slides = [SlideType('default'),SlideType_h5([('SWM',['1','2','3']),('SWM',['1','2','3'])])]
# #
# #     for Sample_sld in slides:
# #         Sample_sld.generate()

# class slideType_definition(SlideType):
#     def __init__(self, defStr):
#         self._defStr = defStr

# class slideType_head_default(SlideType):
#     def __init__(self, h5Tuples):
#         self._h5Tuples = h5Tuples
#         self._lines = h5Tuples[0][0]
# class slideType_head_timeLine(SlideType):
#     def __init__(self, h5Tuples):
#         self._h5Tuples = h5Tuples
#         self._lines = h5Tuples[0][0]
# class slideType_head_multiLine(SlideType):
#     def __init__(self, h5Tuples):
#         self._h5Tuples = h5Tuples
#         self._lines = h5Tuples[0][0]
# class slideType_default(SlideType):
#     def __init__(self, lines):
#         self._lines = lines
# class slideType_singleLine(SlideType):#definition
#     def __init__(self, lines):
#         self._lines = lines
# class slideType_timeLine(SlideType):
#     def __init__(self, lines):
#         self._lines = lines
# class slideType_multiLine(SlideType):
#     def __init__(self, lines):
#         self._lines = lines
class SlideType_title (SlideType):
    def __init__(self, titleTuple, imageLinks):
        self._titleTuple = titleTuple
        self._imageLinks = imageLinks

class SlideType_midTitle (SlideType):
    def __init__(self, midTitle, imageLinks):
        self._midTitle = midTitle
        self._imageLinks = imageLinks

class SlideType_singleLine (SlideType):
    def __init__(self, title, text, imageLinks):
        self._title = title
        self._text = text
        self._imageLinks = imageLinks

class SlideType_multiLine (SlideType):
    def __init__(self, title, textList, imageLinks):
        self._title = title
        self._textList = textList
        self._imageLinks = imageLinks

class SlideType_timeLine (SlideType):
    def __init__(self, title, textList, imageLinks):
        self._title = title
        self._textList = textList
        self._imageLinks = imageLinks

class SlideType_default (SlideType):
    def __init__(self, title, contentsList, imageLinks):
        self._title = title
        self._contentsList = contentsList
        self._imageLinks = imageLinks

class SlideType_head_timeLine (SlideType):
    def __init__(self, title, headTuples, imageLinks):
        self._title = title
        self._headTuples = headTuples
        self._imageLinks = imageLinks

class SlideType_head_multiLine (SlideType):
    def __init__(self, title, headTuples, imageLinks):
        self._title = title
        self._headTuples = headTuples
        self._imageLinks = imageLinks

class SlideType_head_default (SlideType):
    def __init__(self, title, headTuples, imageLinks):
        self._title = title
        self._headTuples = headTuples
        self._imageLinks = imageLinks
