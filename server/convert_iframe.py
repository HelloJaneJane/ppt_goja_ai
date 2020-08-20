#-*- coding:utf-8 -*-
from pptx import Presentation
# import xml.etree.ElementTree as elemTree

def convert_iframe(htmlStr):
    prs = Presentation()#"theme.pptx" 테마설정
    # tree = elemTree.fromstring(htmlStr)

    htmlList = htmlStr.split('<div><br></div>')
    print(htmlList)

    slideCnt = 0
    for element in htmlList:
        lines = element.split('<div>')
        print(lines)
    
        # 제목슬라이드는 layout 0번
        if slideCnt == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            #첫줄은 제목 둘째줄은 부제목
            title.text = lines[0]
            subtitle.text = lines[1].replace('</div>','')

        # 내용슬라이드는 layout 1번
        else :
            if len(lines)<=1:
                continue

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            # [0]:빈칸 / [1]:제목 / [2~]:내용
            title_shape.text = lines[1].replace('</div>','')
            for i in range(2,len(lines)):
                tf = body_shape.text_frame
                p = tf.add_paragraph()
                p.text = lines[i].replace('</div>','')

        slideCnt = slideCnt + 1

    #출력
    prs.save('output.pptx')